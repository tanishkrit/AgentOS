"""
Presentation Agent — Autonomous PowerPoint Generation

Creates professional, styled .pptx presentations using python-pptx.
Leverages the local LLM (Ollama) to generate slide content from
research data or user instructions, then builds the deck with:

  - Widescreen 16:9 layout
  - Dark / Light color themes
  - Accent shapes, divider lines, and styled cards
  - Slide-to-slide fade transitions (raw OpenXML injection)
  - Properly formatted text with font sizing and alignment

All processing is 100% local — no cloud APIs.
"""

import json
import logging
import os
import re
import subprocess
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass, field

from src.agents.base_agent import BaseAgent
from src.utils.llm_client import LLMClient
from src.config import Config

logger = logging.getLogger(__name__)

# ── Lazy import python-pptx (install if missing) ────────────────────

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. Attempting auto-install...")
    try:
        import sys
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "python-pptx"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        PPTX_AVAILABLE = True
        logger.info("python-pptx installed successfully.")
    except Exception as e:
        logger.error(f"Failed to install python-pptx: {e}")


# ═══════════════════════════════════════════════════════════════════
# Color Themes
# ═══════════════════════════════════════════════════════════════════

@dataclass
class SlideTheme:
    """Color palette for a presentation theme."""
    name: str
    bg: str           # Background hex (no #)
    title_color: str   # Title text hex
    body_color: str    # Body text hex
    accent: str        # Primary accent hex
    accent2: str       # Secondary accent hex
    card_bg: str       # Card / shape fill hex
    card_border: str   # Card border hex
    subtitle_color: str  # Subtitle text hex


THEMES: dict[str, SlideTheme] = {
    "charcoal_dark": SlideTheme(
        name="Charcoal Dark",
        bg="1A1A2E",
        title_color="EAEAEA",
        body_color="C8C8D0",
        accent="6C63FF",
        accent2="FF6584",
        card_bg="25253D",
        card_border="3A3A5C",
        subtitle_color="A0A0B8",
    ),
    "deep_navy": SlideTheme(
        name="Deep Navy",
        bg="0B1120",
        title_color="F0F4FF",
        body_color="B0BEC5",
        accent="00B4D8",
        accent2="48CAE4",
        card_bg="131B2E",
        card_border="1E2D45",
        subtitle_color="8899AA",
    ),
    "midnight_violet": SlideTheme(
        name="Midnight Violet",
        bg="110F1A",
        title_color="F5F0FF",
        body_color="C4B5D4",
        accent="8B5CF6",
        accent2="A78BFA",
        card_bg="1A1726",
        card_border="2D2640",
        subtitle_color="9A8BBF",
    ),
    "clean_light": SlideTheme(
        name="Clean Light",
        bg="FAFBFF",
        title_color="1A1A2E",
        body_color="4A4A5A",
        accent="4361EE",
        accent2="3A86FF",
        card_bg="FFFFFF",
        card_border="E2E4EE",
        subtitle_color="6C6C80",
    ),
}


# ═══════════════════════════════════════════════════════════════════
# LLM Prompt for Slide Structure Generation
# ═══════════════════════════════════════════════════════════════════

SLIDE_GENERATION_SYSTEM = """You are a professional presentation designer.
Given a topic and research data, create a structured slide deck outline.

Return ONLY valid JSON in this exact format:
{
  "title": "Presentation Main Title",
  "subtitle": "A compelling subtitle",
  "theme": "charcoal_dark",
  "slides": [
    {
      "type": "section",
      "title": "Section Title",
      "subtitle": "Optional section subtitle"
    },
    {
      "type": "content",
      "title": "Slide Title",
      "bullets": ["Point 1", "Point 2", "Point 3", "Point 4"]
    },
    {
      "type": "two_column",
      "title": "Comparison Title",
      "left_title": "Left Column",
      "left_bullets": ["Point A", "Point B"],
      "right_title": "Right Column",
      "right_bullets": ["Point C", "Point D"]
    },
    {
      "type": "content",
      "title": "Another Slide",
      "bullets": ["Detail 1", "Detail 2", "Detail 3"]
    }
  ]
}

RULES:
- Create 8-15 slides for a comprehensive presentation.
- Available slide types: "section" (for section dividers), "content" (title + bullet points), "two_column" (side-by-side comparison).
- Each "content" slide should have 3-5 bullet points.
- Each bullet point should be a clear, concise sentence (10-20 words max).
- Available themes: "charcoal_dark", "deep_navy", "midnight_violet", "clean_light".
- If the user mentions "dark" choose a dark theme. If they mention "light" choose "clean_light".
- The first slide is always the title slide (auto-generated from "title" and "subtitle").
- The last slide should be a section-type "Thank You" or "Conclusion" slide.
- Make the content informative, well-structured, and professional.
- Do NOT include a title slide in the "slides" array — it is generated automatically.
"""


# ═══════════════════════════════════════════════════════════════════
# Presentation Agent
# ═══════════════════════════════════════════════════════════════════


class PresentationAgent(BaseAgent):
    """
    Generates professional PowerPoint presentations locally.

    Workflow:
    1. Collects research data from upstream dependency agents.
    2. Queries the local LLM to design slide structure and content.
    3. Builds the .pptx file with python-pptx (themes, shapes, transitions).
    4. Saves to the workspace directory.
    """

    def execute(
        self,
        task_description: str,
        parameters: dict,
        dependency_data: dict,
    ) -> dict:
        """
        Execute the presentation creation task.

        Expected parameters:
            - theme (str): Optional theme name override.
            - filename (str): Optional output filename.
            - raw_goal (str): The user's original goal text.
        """
        self.logger.info(f"Starting presentation task: {task_description}")
        self._emit_log(f"📊 Presentation Agent starting: {task_description}")

        if not PPTX_AVAILABLE:
            self._emit_log("❌ python-pptx is not available. Cannot create presentation.")
            return {
                "success": False,
                "summary": "python-pptx library is not installed.",
                "error": "python-pptx library is not installed",
            }

        llm = LLMClient.get_instance()
        if not llm.available:
            self._emit_log("⚠️ LLM unavailable — generating a basic template presentation.")
            return self._generate_fallback_presentation(task_description, parameters)

        # ── Gather research data from upstream agents ────────────────
        research_text = self._collect_research_data(dependency_data)
        raw_goal = parameters.get("raw_goal", task_description)

        # ── Ask LLM to design the slide deck (with retry) ────────────
        slide_data = self._generate_slide_data_with_retry(llm, raw_goal, research_text)

        if not slide_data or "slides" not in slide_data:
            self._emit_log("⚠️ LLM returned invalid slide data after retries. Using fallback.")
            return self._generate_fallback_presentation(task_description, parameters)

        # Sanitize slide data to prevent crashes from malformed LLM output
        slide_data = self._sanitize_slide_data(slide_data)

        self._emit_log(
            f"✅ LLM designed {len(slide_data.get('slides', []))} slides "
            f"with theme '{slide_data.get('theme', 'charcoal_dark')}'"
        )

        # ── Build the .pptx file ─────────────────────────────────────
        theme_override = parameters.get("theme")
        theme_name = theme_override or slide_data.get("theme", "charcoal_dark")
        theme = THEMES.get(theme_name, THEMES["charcoal_dark"])

        # Detect dark/light preference from goal text
        if not theme_override:
            goal_lower = raw_goal.lower()
            if "light" in goal_lower or "white" in goal_lower:
                theme = THEMES["clean_light"]
            elif "navy" in goal_lower or "blue" in goal_lower:
                theme = THEMES["deep_navy"]
            elif "violet" in goal_lower or "purple" in goal_lower:
                theme = THEMES["midnight_violet"]

        self._emit_log(f"🎨 Using theme: {theme.name}")

        try:
            filepath = self._build_presentation(slide_data, theme, parameters)
            self._emit_log(f"✅ Presentation saved: {filepath}")

            # Try to open the file
            self._open_file(filepath)

            return {
                "success": True,
                "summary": f"Created {len(slide_data.get('slides', []))+1}-slide presentation: {filepath}",
                "filepath": filepath,
            }
        except Exception as e:
            self.logger.error(f"Presentation build failed: {e}", exc_info=True)
            self._emit_log(f"❌ Failed to build presentation: {e}")
            # Try fallback as last resort
            self._emit_log("🔄 Attempting fallback presentation...")
            return self._generate_fallback_presentation(task_description, parameters)

    def _generate_slide_data_with_retry(
        self, llm: "LLMClient", raw_goal: str, research_text: str,
    ) -> dict | None:
        """Try to generate slide data from LLM with up to 2 attempts."""
        for attempt in range(2):
            self._emit_log(
                f"🤖 Asking LLM to design slide structure"
                f"{' (retry with shorter context)...' if attempt > 0 else '...'}"
            )

            prompt = f"Topic: {raw_goal}\n\n"
            # On retry, use less research text to reduce LLM confusion
            max_research = 6000 if attempt == 0 else 2000
            if research_text:
                prompt += (
                    f"Research data to use for content:\n"
                    f"\"\"\"\n{research_text[:max_research]}\n\"\"\"\n\n"
                )
            prompt += "Create a professional, informative presentation based on this topic and data."

            try:
                slide_data = llm.generate_json(
                    prompt=prompt,
                    system=SLIDE_GENERATION_SYSTEM,
                    temperature=0.4,
                    max_tokens=4096,
                    timeout=180 if attempt == 0 else 120,
                )

                if slide_data and "slides" in slide_data and len(slide_data["slides"]) > 0:
                    return slide_data

                self.logger.warning(
                    f"LLM attempt {attempt+1} returned invalid data: "
                    f"{list(slide_data.keys()) if slide_data else 'None'}"
                )
            except Exception as e:
                self.logger.warning(f"LLM attempt {attempt+1} failed: {e}")

        return None

    def _sanitize_slide_data(self, slide_data: dict) -> dict:
        """
        Sanitize LLM-generated slide data to prevent crashes.

        Ensures all required fields exist with correct types.
        """
        slide_data.setdefault("title", "Presentation")
        slide_data.setdefault("subtitle", "")
        slide_data.setdefault("theme", "charcoal_dark")

        # Ensure title/subtitle are strings
        if not isinstance(slide_data["title"], str):
            slide_data["title"] = str(slide_data["title"])
        if not isinstance(slide_data["subtitle"], str):
            slide_data["subtitle"] = str(slide_data["subtitle"]) if slide_data["subtitle"] else ""

        clean_slides = []
        for slide in slide_data.get("slides", []):
            if not isinstance(slide, dict):
                continue

            slide_type = slide.get("type", "content")
            if slide_type not in ("section", "content", "two_column"):
                slide_type = "content"
                slide["type"] = slide_type

            # Ensure title is always a string
            slide.setdefault("title", "Untitled Slide")
            if not isinstance(slide["title"], str):
                slide["title"] = str(slide["title"])

            if slide_type == "content":
                bullets = slide.get("bullets", [])
                if not isinstance(bullets, list):
                    bullets = [str(bullets)] if bullets else ["No content"]
                # Ensure all bullets are strings
                slide["bullets"] = [str(b) for b in bullets if b]
                if not slide["bullets"]:
                    slide["bullets"] = ["No content"]

            elif slide_type == "two_column":
                slide.setdefault("left_title", "")
                slide.setdefault("right_title", "")
                left_b = slide.get("left_bullets", [])
                right_b = slide.get("right_bullets", [])
                slide["left_bullets"] = [str(b) for b in left_b if b] if isinstance(left_b, list) else [str(left_b)]
                slide["right_bullets"] = [str(b) for b in right_b if b] if isinstance(right_b, list) else [str(right_b)]

            elif slide_type == "section":
                slide.setdefault("subtitle", "")
                if not isinstance(slide.get("subtitle", ""), str):
                    slide["subtitle"] = str(slide["subtitle"]) if slide["subtitle"] else ""

            clean_slides.append(slide)

        slide_data["slides"] = clean_slides
        return slide_data

    # ─────────────────────────────────────────────────────────────────
    # Research Data Collection
    # ─────────────────────────────────────────────────────────────────

    def _collect_research_data(self, dependency_data: dict) -> str:
        """Extract text content from upstream research agent results."""
        parts = []
        for dep_id, dep_result in dependency_data.items():
            if not isinstance(dep_result, dict):
                continue

            # Structured data from LLM extraction
            structured = dep_result.get("structured_data", [])
            if structured:
                for item in structured[:20]:
                    parts.append(json.dumps(item, default=str))

            # Raw scraped page content
            data_list = dep_result.get("data", [])
            for item in data_list[:10]:
                content = item.get("content", "") or item.get("snippet", "")
                if content:
                    parts.append(content[:1000])

            # Summary
            summary = dep_result.get("summary", "")
            if summary:
                parts.append(summary)

        return "\n\n---\n\n".join(parts)

    # ─────────────────────────────────────────────────────────────────
    # PowerPoint Builder
    # ─────────────────────────────────────────────────────────────────

    def _build_presentation(
        self,
        slide_data: dict,
        theme: SlideTheme,
        parameters: dict,
    ) -> str:
        """Build the actual .pptx file using python-pptx."""
        prs = Presentation()

        # Set widescreen 16:9 dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # ── Title Slide ──────────────────────────────────────────────
        self._add_title_slide(
            prs, theme, slide_w, slide_h,
            title=slide_data.get("title", "Presentation"),
            subtitle=slide_data.get("subtitle", ""),
        )

        # ── Content Slides ───────────────────────────────────────────
        slides = slide_data.get("slides", [])
        slides_built = 0
        for idx, slide_info in enumerate(slides):
            slide_type = slide_info.get("type", "content")
            slide_title = slide_info.get("title", "Untitled")
            self._emit_log(
                f"📝 Building slide {idx+2}/{len(slides)+1}: "
                f"{slide_title[:40]}"
            )

            try:
                if slide_type == "section":
                    self._add_section_slide(
                        prs, theme, slide_w, slide_h,
                        title=slide_info.get("title", "Section"),
                        subtitle=slide_info.get("subtitle", ""),
                    )
                elif slide_type == "two_column":
                    self._add_two_column_slide(
                        prs, theme, slide_w, slide_h,
                        title=slide_info.get("title", ""),
                        left_title=slide_info.get("left_title", ""),
                        left_bullets=slide_info.get("left_bullets", []),
                        right_title=slide_info.get("right_title", ""),
                        right_bullets=slide_info.get("right_bullets", []),
                    )
                else:  # "content"
                    self._add_content_slide(
                        prs, theme, slide_w, slide_h,
                        title=slide_info.get("title", ""),
                        bullets=slide_info.get("bullets", []),
                    )
                slides_built += 1
            except Exception as e:
                # Log and skip broken slides instead of crashing
                self.logger.warning(
                    f"Skipped slide {idx+2} ({slide_title[:30]}): {e}"
                )
                self._emit_log(
                    f"⚠️ Skipped slide '{slide_title[:30]}' due to error: {e}"
                )

        if slides_built == 0:
            raise RuntimeError("All slides failed to build — presentation is empty.")

        # ── Add fade transitions to all slides ──────────────────────
        try:
            self._inject_transitions(prs)
        except Exception as e:
            self.logger.warning(f"Transition injection failed (non-fatal): {e}")

        # ── Save ─────────────────────────────────────────────────────
        filename = parameters.get("filename")
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = re.sub(r'[^\w\s-]', '', slide_data.get("title", "presentation"))
            safe_title = re.sub(r'\s+', '_', safe_title.strip())[:40]
            filename = f"{safe_title}_{timestamp}.pptx"
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        output_path = Config.WORKSPACE_ROOT / filename
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        return str(output_path)

    # ─────────────────────────────────────────────────────────────────
    # Slide Builders
    # ─────────────────────────────────────────────────────────────────

    def _add_title_slide(
        self, prs, theme: SlideTheme, sw, sh,
        title: str, subtitle: str,
    ) -> None:
        """Add a styled title slide with accent line and background."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._set_bg(slide, theme.bg)

        # Accent bar at top
        self._add_shape(
            slide, left=0, top=0,
            width=sw, height=Inches(0.08),
            fill=theme.accent,
        )

        # Decorative accent circle
        self._add_shape(
            slide,
            left=sw - Inches(2.5), top=Inches(0.6),
            width=Inches(1.8), height=Inches(1.8),
            fill=theme.accent, opacity=15,
            shape_type="oval",
        )

        # Title text
        self._add_textbox(
            slide,
            left=Inches(1.2), top=Inches(2.0),
            width=sw - Inches(2.4), height=Inches(2.0),
            text=title,
            font_size=44, bold=True,
            color=theme.title_color,
            alignment=PP_ALIGN.LEFT,
        )

        # Accent divider line
        self._add_shape(
            slide,
            left=Inches(1.2), top=Inches(4.1),
            width=Inches(3.0), height=Inches(0.06),
            fill=theme.accent,
        )

        # Subtitle
        if subtitle:
            self._add_textbox(
                slide,
                left=Inches(1.2), top=Inches(4.5),
                width=sw - Inches(2.4), height=Inches(1.2),
                text=subtitle,
                font_size=20,
                color=theme.subtitle_color,
                alignment=PP_ALIGN.LEFT,
            )

        # Date
        self._add_textbox(
            slide,
            left=Inches(1.2), top=sh - Inches(1.0),
            width=Inches(4), height=Inches(0.5),
            text=datetime.now().strftime("%B %d, %Y"),
            font_size=12,
            color=theme.subtitle_color,
            alignment=PP_ALIGN.LEFT,
        )

    def _add_section_slide(
        self, prs, theme: SlideTheme, sw, sh,
        title: str, subtitle: str = "",
    ) -> None:
        """Add a section divider slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide, theme.bg)

        # Large accent circle (decorative)
        self._add_shape(
            slide,
            left=Inches(-1), top=Inches(-1),
            width=Inches(4), height=Inches(4),
            fill=theme.accent, opacity=10,
            shape_type="oval",
        )

        # Accent line
        self._add_shape(
            slide,
            left=Inches(1.5), top=Inches(3.0),
            width=Inches(2.5), height=Inches(0.06),
            fill=theme.accent,
        )

        # Title
        self._add_textbox(
            slide,
            left=Inches(1.5), top=Inches(3.3),
            width=sw - Inches(3.0), height=Inches(1.8),
            text=title,
            font_size=38, bold=True,
            color=theme.title_color,
            alignment=PP_ALIGN.LEFT,
        )

        # Subtitle
        if subtitle:
            self._add_textbox(
                slide,
                left=Inches(1.5), top=Inches(4.8),
                width=sw - Inches(3.0), height=Inches(1.0),
                text=subtitle,
                font_size=18,
                color=theme.subtitle_color,
                alignment=PP_ALIGN.LEFT,
            )

    def _add_content_slide(
        self, prs, theme: SlideTheme, sw, sh,
        title: str, bullets: list[str],
    ) -> None:
        """Add a content slide with title and bullet points."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide, theme.bg)

        # Accent bar at top
        self._add_shape(
            slide, left=0, top=0,
            width=sw, height=Inches(0.05),
            fill=theme.accent,
        )

        # Title
        self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(0.5),
            width=sw - Inches(2.0), height=Inches(1.0),
            text=title,
            font_size=30, bold=True,
            color=theme.title_color,
            alignment=PP_ALIGN.LEFT,
        )

        # Divider line under title
        self._add_shape(
            slide,
            left=Inches(1.0), top=Inches(1.45),
            width=Inches(2.0), height=Inches(0.04),
            fill=theme.accent,
        )

        # Bullet points as styled cards
        bullet_top = Inches(1.8)
        bullet_spacing = Inches(0.95)

        for i, bullet in enumerate(bullets[:6]):
            card_top = bullet_top + i * bullet_spacing

            # Card background
            self._add_shape(
                slide,
                left=Inches(1.0), top=card_top,
                width=sw - Inches(2.0), height=Inches(0.8),
                fill=theme.card_bg,
                border_color=theme.card_border,
                corner_radius=8,
            )

            # Accent dot
            self._add_shape(
                slide,
                left=Inches(1.3), top=card_top + Inches(0.28),
                width=Inches(0.2), height=Inches(0.2),
                fill=theme.accent,
                shape_type="oval",
            )

            # Bullet text
            self._add_textbox(
                slide,
                left=Inches(1.7), top=card_top + Inches(0.12),
                width=sw - Inches(3.0), height=Inches(0.55),
                text=bullet,
                font_size=16,
                color=theme.body_color,
                alignment=PP_ALIGN.LEFT,
            )

    def _add_two_column_slide(
        self, prs, theme: SlideTheme, sw, sh,
        title: str,
        left_title: str, left_bullets: list[str],
        right_title: str, right_bullets: list[str],
    ) -> None:
        """Add a two-column comparison slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide, theme.bg)

        # Accent bar at top
        self._add_shape(
            slide, left=0, top=0,
            width=sw, height=Inches(0.05),
            fill=theme.accent,
        )

        # Title
        self._add_textbox(
            slide,
            left=Inches(1.0), top=Inches(0.5),
            width=sw - Inches(2.0), height=Inches(1.0),
            text=title,
            font_size=30, bold=True,
            color=theme.title_color,
            alignment=PP_ALIGN.LEFT,
        )

        # Divider
        self._add_shape(
            slide,
            left=Inches(1.0), top=Inches(1.45),
            width=Inches(2.0), height=Inches(0.04),
            fill=theme.accent,
        )

        col_width = (sw - Inches(2.5)) // 2

        # Left column card
        self._add_shape(
            slide,
            left=Inches(1.0), top=Inches(1.8),
            width=col_width, height=sh - Inches(2.5),
            fill=theme.card_bg,
            border_color=theme.card_border,
            corner_radius=10,
        )

        # Left column title
        self._add_textbox(
            slide,
            left=Inches(1.3), top=Inches(2.0),
            width=col_width - Inches(0.6), height=Inches(0.6),
            text=left_title,
            font_size=20, bold=True,
            color=theme.accent,
            alignment=PP_ALIGN.LEFT,
        )

        # Left column bullets
        for i, bullet in enumerate(left_bullets[:5]):
            self._add_textbox(
                slide,
                left=Inches(1.5), top=Inches(2.7) + i * Inches(0.7),
                width=col_width - Inches(0.8), height=Inches(0.6),
                text=f"• {bullet}",
                font_size=14,
                color=theme.body_color,
                alignment=PP_ALIGN.LEFT,
            )

        # Right column card
        right_left = Inches(1.0) + col_width + Inches(0.5)
        self._add_shape(
            slide,
            left=right_left, top=Inches(1.8),
            width=col_width, height=sh - Inches(2.5),
            fill=theme.card_bg,
            border_color=theme.card_border,
            corner_radius=10,
        )

        # Right column title
        self._add_textbox(
            slide,
            left=right_left + Inches(0.3), top=Inches(2.0),
            width=col_width - Inches(0.6), height=Inches(0.6),
            text=right_title,
            font_size=20, bold=True,
            color=theme.accent2,
            alignment=PP_ALIGN.LEFT,
        )

        # Right column bullets
        for i, bullet in enumerate(right_bullets[:5]):
            self._add_textbox(
                slide,
                left=right_left + Inches(0.5), top=Inches(2.7) + i * Inches(0.7),
                width=col_width - Inches(0.8), height=Inches(0.6),
                text=f"• {bullet}",
                font_size=14,
                color=theme.body_color,
                alignment=PP_ALIGN.LEFT,
            )

    # ─────────────────────────────────────────────────────────────────
    # Low-Level Helpers
    # ─────────────────────────────────────────────────────────────────

    def _set_bg(self, slide, hex_color: str) -> None:
        """Set a solid background color on a slide."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(hex_color)

    def _add_textbox(
        self, slide,
        left, top, width, height,
        text: str,
        font_size: int = 16,
        bold: bool = False,
        color: str = "FFFFFF",
        alignment=PP_ALIGN.LEFT,
    ) -> None:
        """Add a styled text box to a slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor.from_string(color)
        p.alignment = alignment

        # Use Segoe UI for a modern look
        p.font.name = "Segoe UI"

    def _add_shape(
        self, slide,
        left, top, width, height,
        fill: str = "FFFFFF",
        border_color: str | None = None,
        opacity: int = 100,
        shape_type: str = "rect",
        corner_radius: int = 0,
    ) -> None:
        """Add a shape (rectangle or oval) with fill and optional border."""
        from pptx.enum.shapes import MSO_SHAPE

        if shape_type == "oval":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, top, width, height
            )
        elif corner_radius > 0:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
            )
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )

        # Fill
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill)

        # Opacity (via alpha on the fill)
        if opacity < 100:
            # python-pptx doesn't directly support opacity, so we modify XML
            solidFill = shape.fill._fill._solidFill
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is None:
                srgb = solidFill.find(qn('a:schemeClr'))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha.set('val', str(opacity * 1000))  # value in 1000ths of a percent

        # Border
        if border_color:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(border_color)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()  # No border

    def _inject_transitions(self, prs) -> None:
        """
        Inject slide-to-slide fade transitions via raw OpenXML.

        python-pptx doesn't support transitions natively, so we inject
        the <p:transition> element directly into each slide's XML.
        """
        nsmap = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        }

        for slide in prs.slides:
            slide_elem = slide._element
            # Remove any existing transition
            for existing in slide_elem.findall(qn('p:transition')):
                slide_elem.remove(existing)

            # Create <p:transition spd="med" advClick="1">
            #   <p:fade />
            # </p:transition>
            transition = etree.SubElement(slide_elem, qn('p:transition'))
            transition.set('spd', 'med')
            transition.set('advClick', '1')
            fade = etree.SubElement(transition, qn('p:fade'))

    # ─────────────────────────────────────────────────────────────────
    # Fallback (no LLM)
    # ─────────────────────────────────────────────────────────────────

    def _generate_fallback_presentation(
        self, task_description: str, parameters: dict,
    ) -> dict:
        """Generate a minimal presentation when LLM is unavailable."""
        self._emit_log("📊 Generating fallback presentation...")

        fallback_data = {
            "title": task_description[:60],
            "subtitle": "Auto-generated presentation",
            "theme": "charcoal_dark",
            "slides": [
                {
                    "type": "section",
                    "title": "Overview",
                    "subtitle": task_description,
                },
                {
                    "type": "content",
                    "title": "Key Points",
                    "bullets": [
                        "Research data was collected from web sources",
                        "Content generated by local AI agent",
                        "Presentation created autonomously",
                        "Edit slides as needed for your use case",
                    ],
                },
                {
                    "type": "section",
                    "title": "Thank You",
                    "subtitle": "Generated by AgentOS",
                },
            ],
        }

        theme = THEMES["charcoal_dark"]
        try:
            filepath = self._build_presentation(fallback_data, theme, parameters)
            self._emit_log(f"✅ Fallback presentation saved: {filepath}")
            self._open_file(filepath)
            return {
                "success": True,
                "summary": f"Created fallback presentation: {filepath}",
                "filepath": filepath,
            }
        except Exception as e:
            return {
                "success": False,
                "summary": f"Fallback presentation failed: {e}",
            }

    def _open_file(self, filepath: str) -> None:
        """Attempt to open the generated file with the default application."""
        try:
            os.startfile(filepath)
            self._emit_log(f"📂 Opened: {filepath}")
        except Exception:
            self._emit_log(f"📂 File saved at: {filepath} (open manually)")
