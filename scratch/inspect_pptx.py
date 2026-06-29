import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, 0, 0, 100000, 100000)

print("Checking shape.line properties...")
print("shape.line type:", type(shape.line))
print("shape.line attributes:", dir(shape.line))

try:
    print("Trying shape.line.fill...")
    print("shape.line.fill type:", type(shape.line.fill))
    print("shape.line.fill attributes:", dir(shape.line.fill))
except Exception as e:
    print("shape.line.fill failed:", e)

try:
    print("Trying shape.line.fill.background()...")
    shape.line.fill.background()
    print("shape.line.fill.background() Succeeded!")
except Exception as e:
    print("shape.line.fill.background() Failed:", e)
